import streamlit as st
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from PIL import Image
import io

# Slide size options
slide_size_map = {
    "4:3": (10, 7.5),
    "16:9": (13.33, 7.5),
    "Widescreen": (16, 9)
}

# Placement positions
def get_position(placement, img_w, img_h, slide_w, slide_h):
    positions = {
        "Top Left": (0.5, 0.5),
        "Top Right": (slide_w - img_w - 0.5, 0.5),
        "Bottom Left": (0.5, slide_h - img_h - 0.5),
        "Bottom Right": (slide_w - img_w - 0.5, slide_h - img_h - 0.5),
        "Center": ((slide_w - img_w) / 2, (slide_h - img_h) / 2),
        "Top Middle": ((slide_w - img_w) / 2, 0.5),
        "Bottom Middle": ((slide_w - img_w) / 2, slide_h - img_h - 0.5)
    }
    return positions[placement]

# Streamlit GUI
st.title("PDF to PPT Converter with Margin Options")

# Upload multiple PDFs
pdf_files = st.file_uploader("Upload one or more PDFs to merge", type=["pdf"], accept_multiple_files=True)

# Background color options
bg_color = st.selectbox("Background Color", ["White", "LightGray", "LightBlue", "Yellow"])
color_map = {
    "White": "FFFFFF",
    "LightGray": "D3D3D3",
    "LightBlue": "ADD8E6",
    "Yellow": "FFFF99"
}

# Shrink % and placement
shrink = st.slider("Shrink PDF Page (%)", 25, 100, 50) / 100
placement = st.selectbox("Placement", ["Top Left", "Top Right", "Bottom Left", "Bottom Right", "Center", "Top Middle", "Bottom Middle"])
slide_size_label = st.selectbox("Slide Size", ["4:3", "16:9", "Widescreen"])

# Generate PPT on button click
if pdf_files and st.button("Generate Merged PPT"):
    ppt = Presentation()
    slide_w, slide_h = slide_size_map[slide_size_label]
    ppt.slide_width = Inches(slide_w)
    ppt.slide_height = Inches(slide_h)

    for uploaded_pdf in pdf_files:
        doc = fitz.open(stream=uploaded_pdf.read(), filetype="pdf")
        for page in doc:
            pix = page.get_pixmap(dpi=150)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            img_w_in = (pix.width / 150) * shrink
            img_h_in = (pix.height / 150) * shrink
            img_resized = img.resize((int(pix.width * shrink), int(pix.height * shrink)))
            img_stream = io.BytesIO()
            img_resized.save(img_stream, format='PNG')
            img_stream.seek(0)

            slide = ppt.slides.add_slide(ppt.slide_layouts[6])
            slide.background.fill.solid()
            hex_color = color_map[bg_color]
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
            slide.background.fill.fore_color.rgb = RGBColor(r, g, b)

            left, top = get_position(placement, img_w_in, img_h_in, slide_w, slide_h)
            slide.shapes.add_picture(img_stream, Inches(left), Inches(top), width=Inches(img_w_in), height=Inches(img_h_in))

    ppt_stream = io.BytesIO()
    ppt.save(ppt_stream)
    ppt_stream.seek(0)
    st.download_button("Download Merged PPT", ppt_stream, file_name="merged_ppt.pptx")
